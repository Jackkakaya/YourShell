import time, os
R = os.environ['HOME'] + '/Documents/pyperf.txt'
open(R,'w').write('=== Python 3.14 生成性能 (iOS sim) ===\n')
def w(l): open(R,'a').write(l+'\n')

# 1. PPTX: 20-slide deck with titles + content
t=time.time()
from pptx import Presentation
from pptx.util import Inches, Pt
prs = Presentation()
for i in range(20):
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = f'Slide {i+1}'
    s.placeholders[1].text = '\n'.join(f'bullet point {j}' for j in range(5))
prs.save(os.environ['HOME']+'/Documents/out.pptx')
sz = os.path.getsize(os.environ['HOME']+'/Documents/out.pptx')
w(f'PPTX 20 slides:          {int((time.time()-t)*1000)} ms ({sz//1024}KB)')

# 2. PIL: draw a bar chart 800x600
t=time.time()
from PIL import Image, ImageDraw, ImageFont
img = Image.new('RGB',(800,600),'white')
d = ImageDraw.Draw(img)
data=[120,240,180,300,90,270,150,330]
for i,v in enumerate(data):
    x=60+i*90
    d.rectangle([x,560-v,x+60,560],fill=(60,120,220))
    d.text((x+10,565),f'B{i}',fill='black')
d.text((250,20),'YourShell Bar Chart',fill='black')
img.save(os.environ['HOME']+'/Documents/chart.png')
w(f'PIL 800x600 bar chart:   {int((time.time()-t)*1000)} ms')

# 3. PIL heavy: 2000x2000 gradient pixel-by-pixel (compute heavy)
t=time.time()
img2=Image.new('RGB',(500,500))
px=img2.load()
for y in range(500):
    for x in range(500):
        px[x,y]=(x%256,y%256,(x+y)%256)
img2.save(os.environ['HOME']+'/Documents/grad.png')
w(f'PIL 500x500 pixel loop:  {int((time.time()-t)*1000)} ms')

# 4. PDF: 30-page document with text + the chart image
t=time.time()
from fpdf import FPDF
pdf=FPDF()
for i in range(30):
    pdf.add_page()
    pdf.set_font('Helvetica',size=16)
    pdf.cell(text=f'Page {i+1} of the YourShell report')
    pdf.ln(20)
    pdf.set_font('Helvetica',size=10)
    for j in range(10):
        pdf.cell(text=f'  line {j}: some report content here'); pdf.ln(6)
pdf.image(os.environ['HOME']+'/Documents/chart.png',x=10,y=200,w=80)
pdf.output(os.environ['HOME']+'/Documents/report.pdf')
sz=os.path.getsize(os.environ['HOME']+'/Documents/report.pdf')
w(f'PDF 30 pages + image:    {int((time.time()-t)*1000)} ms ({sz//1024}KB)')
w('DONE')
